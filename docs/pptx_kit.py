# -*- coding: utf-8 -*-
"""arch-pptx 공용 렌더링 키트.

아키텍처 분석 PPTX를 일관된 스타일(딥 네이비 + 브랜드 블루)로 그리는 헬퍼 모음.
gen_pptx.py 가 `from pptx_kit import *` 로 가져다 쓴다.

저수준 프리미티브(box/txt/arrow/connector)와 고수준 슬라이드 빌더(cover/
bullet_cards/section_table/revision_history/summary)를 함께 제공한다.
텍스트·표 중심 슬라이드는 고수준 빌더로, 다이어그램은 프리미티브로 직접 그린다.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 팔레트 ──────────────────────────────────────────────
NAVY   = RGBColor(0x0F, 0x1B, 0x33)
BRAND  = RGBColor(0x2F, 0x5B, 0xEA)
CYAN   = RGBColor(0x0C, 0x97, 0xB8)
PURPLE = RGBColor(0x7A, 0x5A, 0xF8)
GREEN  = RGBColor(0x14, 0x9A, 0x6B)
ORANGE = RGBColor(0xE0, 0x51, 0x2E)
INK    = RGBColor(0x1A, 0x23, 0x33)
GREY   = RGBColor(0x5B, 0x66, 0x77)
LIGHT  = RGBColor(0xF2, 0xF5, 0xFA)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0xD3, 0xDA, 0xE6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

KOR  = "맑은 고딕"
MONO = "Consolas"

# 슬라이드 푸터 좌측 라벨 — 프로젝트별로 set_footer()로 바꾼다.
FOOTER_LABEL = "아키텍처 분석 보고서"


def set_footer(label):
    global FOOTER_LABEL
    FOOTER_LABEL = label


# ── 덱/슬라이드 ──────────────────────────────────────────
def new_deck(width_in=13.333, height_in=7.5):
    prs = Presentation()
    prs.slide_width  = Inches(width_in)
    prs.slide_height = Inches(height_in)
    return prs


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def set_bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


# ── 프리미티브 ───────────────────────────────────────────
def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {})
        sh = ef.makeelement(qn('a:outerShdw'), {'blurRad': '60000', 'dist': '25000', 'dir': '5400000', 'rotWithShape': '0'})
        clr = sh.makeelement(qn('a:srgbClr'), {'val': '1A2333'})
        alpha = clr.makeelement(qn('a:alpha'), {'val': '22000'})
        clr.append(alpha); sh.append(clr); ef.append(sh); el.append(ef)
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, sp_after=2):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    if isinstance(runs, str):
        runs = [(runs, 14, INK, False, KOR)]
    first = True
    for item in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(sp_after)
        p.space_before = Pt(0)
        segs = item if isinstance(item, list) else [item]
        for (t, sz, col, bold, fnt) in segs:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = fnt
    return tb


def label_in(sp, text, size=13, color=WHITE, bold=True, font=KOR, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        segs = ln if isinstance(ln, list) else [(ln, size, color, bold, font)]
        for (t, sz, col, bd, fnt) in segs:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bd; r.font.name = fnt


def connector(s, x1, y1, x2, y2, color=GREY, w=1.75, dash=None):
    cn = s.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    if dash:
        ln = cn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': dash}))
    return cn


def arrow(s, x1, y1, x2, y2, color=GREY, w=1.75, dash=None):
    cn = connector(s, x1, y1, x2, y2, color, w, dash)
    ln = cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return cn


def header(s, kicker, title, accent=BRAND):
    box(s, 0, 0, 13.333, 1.15, fill=NAVY)
    box(s, 0, 1.15, 13.333, 0.06, fill=accent)
    box(s, 0.55, 0.30, 0.13, 0.58, fill=accent)
    txt(s, 0.85, 0.20, 11.5, 0.4, [(kicker, 11, RGBColor(0x8F,0xA6,0xCF), True, KOR)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 0.85, 0.48, 11.8, 0.6, [(title, 23, WHITE, True, KOR)], anchor=MSO_ANCHOR.MIDDLE)


def footer(s, n, label=None):
    txt(s, 0.85, 7.08, 8, 0.3, [(label or FOOTER_LABEL, 8.5, GREY, False, KOR)])
    txt(s, 11.6, 7.08, 1.1, 0.3, [(f"{n:02d}", 8.5, GREY, True, KOR)], align=PP_ALIGN.RIGHT)


def chip(s, x, y, w, text, color):
    h = 0.34
    box(s, x, y, 0.09, h, fill=color)
    b = box(s, x+0.09, y, w-0.09, h, fill=RGBColor(0xF6,0xF8,0xFC), line=LINE, line_w=0.75)
    label_in(b, text, size=10.5, color=INK, bold=True, align=PP_ALIGN.LEFT)
    b.text_frame.margin_left = Pt(8)


# ── 고수준 슬라이드 빌더 ─────────────────────────────────
def cover(prs, kicker, title, subtitle, comps, version_date):
    """표지. comps = [(name, ko, sub, color), ...] (최대 3)."""
    s = add_slide(prs); set_bg(s, NAVY)
    box(s, 0, 0, 13.333, 7.5, fill=NAVY)
    box(s, 9.3, -1.2, 5, 5, fill=RGBColor(0x16,0x29,0x4D), shape=MSO_SHAPE.OVAL)
    box(s, 10.6, 3.6, 4, 4, fill=RGBColor(0x13,0x22,0x40), shape=MSO_SHAPE.OVAL)
    box(s, 0.9, 1.05, 0.16, 1.5, fill=BRAND)
    txt(s, 1.25, 1.0, 10, 0.5, [(kicker, 13, CYAN, True, KOR)])
    txt(s, 1.2, 1.55, 11, 1.6, [
        [(title, 44, WHITE, True, KOR)],
        [(subtitle, 22, RGBColor(0xC9,0xD6,0xEE), True, KOR)],
    ])
    cx = 1.2
    for nm, ko, sub, col in comps[:3]:
        box(s, cx, 3.95, 3.6, 1.55, fill=RGBColor(0x17,0x2A,0x4E), line=col, line_w=1.5, shadow=True)
        box(s, cx, 3.95, 3.6, 0.10, fill=col)
        txt(s, cx+0.28, 4.18, 3.1, 0.4, [(nm, 15, WHITE, True, MONO)])
        txt(s, cx+0.28, 4.62, 3.1, 0.4, [(ko, 15, RGBColor(0x49,0xC4,0xE0) if col == CYAN else col, True, KOR)])
        txt(s, cx+0.28, 5.02, 3.1, 0.4, [(sub, 11, RGBColor(0x9F,0xB2,0xD4), False, KOR)])
        cx += 3.85
    txt(s, 1.2, 6.82, 11, 0.4, [(version_date, 11, GREY, False, KOR)])
    return s


def bullet_cards(prs, n, kicker, title, intro, cards, accent=BRAND):
    """카드 그리드(각 카드 = 제목 + 불릿). cards = [(name, ko, color, [items]), ...] (최대 3)."""
    s = add_slide(prs); set_bg(s, LIGHT)
    header(s, kicker, title, accent=accent)
    if intro:
        txt(s, 0.85, 1.45, 11.6, 0.6, [(intro, 13, GREY, False, KOR)])
    cx = 0.85
    for nm, ko, col, items in cards[:3]:
        box(s, cx, 2.35, 3.82, 4.35, fill=CARD, line=LINE, line_w=1, shadow=True)
        box(s, cx, 2.35, 3.82, 0.62, fill=col)
        txt(s, cx+0.3, 2.45, 3.3, 0.45, [(nm, 15, WHITE, True, MONO)], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, cx+0.3, 3.08, 3.3, 0.4, [(ko, 16, INK, True, KOR)])
        yy = 3.62
        for it in items:
            box(s, cx+0.3, yy+0.07, 0.08, 0.08, fill=col, shape=MSO_SHAPE.OVAL)
            txt(s, cx+0.52, yy-0.04, 3.1, 0.55, [(it, 11, GREY, False, KOR)])
            yy += 0.55
        cx += 4.02
    footer(s, n)
    return s


def section_table(prs, n, kicker, title, columns, rows, intro=None, accent=BRAND, row_h=0.46):
    """단일 표 슬라이드. columns = [(헤더, x, w), ...]; rows = [(셀1, 셀2, ...), ...].
    마지막 컬럼은 보통 설명(넓게). 행은 zebra 배경."""
    s = add_slide(prs); set_bg(s, LIGHT)
    header(s, kicker, title, accent=accent)
    ty = 1.55
    if intro:
        txt(s, 0.85, 1.4, 11.6, 0.4, [(intro, 12, GREY, False, KOR)])
        ty = 1.95
    x0 = columns[0][1]
    total_w = columns[-1][1] + columns[-1][2] - x0
    box(s, x0, ty, total_w, row_h, fill=INK)
    for h_, x, w in columns:
        txt(s, x+0.12, ty+0.07, w-0.2, 0.32, [(h_, 11.5, WHITE, True, KOR)], anchor=MSO_ANCHOR.MIDDLE)
    ty += row_h
    for i, cells in enumerate(rows):
        bg = CARD if i % 2 == 0 else RGBColor(0xEE,0xF2,0xFA)
        box(s, x0, ty, total_w, row_h, fill=bg, line=LINE, line_w=0.5)
        for (h_, x, w), cell in zip(columns, cells):
            mono = isinstance(cell, tuple)  # (text, True) → monospace
            text = cell[0] if mono else cell
            txt(s, x+0.12, ty+0.07, w-0.2, row_h-0.1,
                [(text, 10.5, INK, False, MONO if mono else KOR)], anchor=MSO_ANCHOR.MIDDLE)
        ty += row_h
    footer(s, n)
    return s


def revision_history(prs, n, rows, kicker="REVISION HISTORY", title="문서 변경 이력",
                     intro="프로젝트 버전별 변경 요약."):
    """변경 이력 표. rows = [(버전, 날짜, 유형, 변경내용), ...]. 유형: init/feat/fix."""
    s = add_slide(prs); set_bg(s, LIGHT)
    header(s, kicker, title, accent=GREEN)
    txt(s, 0.85, 1.4, 11.6, 0.4, [(intro, 12, GREY, False, KOR)])
    cols = [("버전", 0.85, 1.45), ("날짜", 2.30, 1.55), ("유형", 3.85, 1.10), ("변경 내용", 4.95, 7.65)]
    ty = 1.95
    box(s, 0.85, ty, 11.75, 0.46, fill=INK)
    for h_, x, w in cols:
        txt(s, x+0.12, ty+0.07, w-0.2, 0.32, [(h_, 11.5, WHITE, True, KOR)], anchor=MSO_ANCHOR.MIDDLE)
    ty += 0.46
    type_col = {"init": BRAND, "feat": GREEN, "fix": ORANGE}
    rh = 0.585
    for i, (ver, dt, typ, desc) in enumerate(rows):
        last = (i == len(rows)-1)
        bg = RGBColor(0xE8,0xF4,0xEE) if last else (CARD if i % 2 == 0 else RGBColor(0xEE,0xF2,0xFA))
        box(s, 0.85, ty, 11.75, rh, fill=bg, line=LINE, line_w=0.5)
        txt(s, cols[0][1]+0.12, ty+0.07, cols[0][2]-0.2, rh-0.1, [(ver, 11.5, INK, True, MONO)], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, cols[1][1]+0.12, ty+0.07, cols[1][2]-0.2, rh-0.1, [(dt, 10.5, GREY, False, MONO)], anchor=MSO_ANCHOR.MIDDLE)
        tb = box(s, cols[2][1]+0.12, ty+0.13, 0.82, rh-0.26, fill=type_col.get(typ, GREY))
        label_in(tb, typ, 10, WHITE, True)
        txt(s, cols[3][1]+0.12, ty+0.07, cols[3][2]-0.2, rh-0.1, [(desc, 10.5, INK if last else GREY, last, KOR)], anchor=MSO_ANCHOR.MIDDLE)
        ty += rh
    footer(s, n)
    return s


def summary(prs, kicker, title, points, version_date):
    """네이비 요약 슬라이드. points = [(제목, 설명, color), ...]."""
    s = add_slide(prs); set_bg(s, NAVY)
    box(s, 0, 0, 13.333, 7.5, fill=NAVY)
    box(s, 9.5, -1.5, 5.5, 5.5, fill=RGBColor(0x15,0x28,0x4B), shape=MSO_SHAPE.OVAL)
    box(s, 0.9, 0.75, 0.16, 1.0, fill=BRAND)
    txt(s, 1.25, 0.7, 10, 0.4, [(kicker, 13, CYAN, True, KOR)])
    txt(s, 1.2, 1.15, 11, 0.7, [(title, 32, WHITE, True, KOR)])
    yy = 2.1
    for t_, d, col in points:
        box(s, 1.2, yy, 11.0, 0.72, fill=RGBColor(0x16,0x29,0x4C), line=col, line_w=1.2)
        box(s, 1.2, yy, 0.12, 0.72, fill=col)
        txt(s, 1.55, yy+0.07, 3.0, 0.55, [(t_, 14, WHITE, True, KOR)], anchor=MSO_ANCHOR.MIDDLE)
        txt(s, 4.7, yy+0.07, 7.3, 0.6, [(d, 11.5, RGBColor(0xC9,0xD6,0xEE), False, KOR)], anchor=MSO_ANCHOR.MIDDLE)
        yy += 0.80
    txt(s, 1.2, 6.95, 11, 0.4, [(version_date, 10, GREY, False, KOR)])
    return s


def save(prs, out_path):
    prs.save(out_path)
    print("SAVED:", out_path)
    print("slides:", len(prs.slides._sldIdLst))
